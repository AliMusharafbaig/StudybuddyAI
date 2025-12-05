"""
StudyBuddy AI - Content Ingestion Agent
=========================================
Processes uploaded materials (PDF, DOCX, video, audio).
"""

import os
import asyncio
from typing import Dict, List, Any, Optional
from pathlib import Path
import logging

from agents.base_agent import BaseAgent
from core.config import settings

logger = logging.getLogger(__name__)


class ContentIngestionAgent(BaseAgent):
    """
    Agent for processing and ingesting course materials.
    
    Supports:
    - PDF extraction (PyPDF2, pdfplumber)
    - DOCX extraction
    - PPTX extraction
    - Audio transcription (Whisper)
    - Video processing (extract audio → transcribe)
    - OCR for images
    """
    
    def __init__(self):
        super().__init__("ContentIngestionAgent")
    
    async def run(
        self,
        file_path: str,
        file_type: str,
        material_id: str = None
    ) -> Dict[str, Any]:
        """
        Process a file and extract text content.
        
        Args:
            file_path: Path to the file
            file_type: File extension (pdf, docx, etc.)
            material_id: ID of the material record
        
        Returns:
            Dict with extracted text, metadata, and chunks
        """
        self._log_action("Processing file", {"file": file_path, "type": file_type})
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        # Extract based on file type
        file_type = file_type.lower()
        
        if file_type == "pdf":
            result = await self._process_pdf(file_path)
        elif file_type in ["docx", "doc"]:
            result = await self._process_docx(file_path)
        elif file_type in ["pptx", "ppt"]:
            result = await self._process_pptx(file_path)
        elif file_type in ["mp3", "wav", "m4a", "ogg"]:
            result = await self._process_audio(file_path)
        elif file_type in ["mp4", "avi", "mov", "mkv"]:
            result = await self._process_video(file_path)
        elif file_type in ["png", "jpg", "jpeg"]:
            result = await self._process_image(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        # Chunk the content
        if result.get("text"):
            result["chunks"] = self._chunk_text(
                result["text"],
                source=os.path.basename(file_path),
                material_id=material_id
            )
        
        self._log_action("File processed", {
            "file": file_path,
            "text_length": len(result.get("text", "")),
            "num_chunks": len(result.get("chunks", []))
        })
        
        return result
    
    async def _process_pdf(self, file_path: str) -> Dict[str, Any]:
        """Extract text from PDF."""
        try:
            import pdfplumber
            
            text_parts = []
            num_pages = 0
            
            with pdfplumber.open(file_path) as pdf:
                num_pages = len(pdf.pages)
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
            
            return {
                "text": "\n\n".join(text_parts),
                "num_pages": num_pages,
                "source_type": "pdf"
            }
        except Exception as e:
            logger.warning(f"pdfplumber failed, trying PyPDF2: {e}")
            
            try:
                from PyPDF2 import PdfReader
                
                reader = PdfReader(file_path)
                text_parts = []
                
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)
                
                return {
                    "text": "\n\n".join(text_parts),
                    "num_pages": len(reader.pages),
                    "source_type": "pdf"
                }
            except Exception as e2:
                logger.error(f"PDF extraction failed: {e2}")
                return {"text": "", "error": str(e2)}
    
    async def _process_docx(self, file_path: str) -> Dict[str, Any]:
        """Extract text from DOCX."""
        try:
            from docx import Document
            
            doc = Document(file_path)
            text_parts = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_parts.append(cell.text)
            
            return {
                "text": "\n\n".join(text_parts),
                "source_type": "docx"
            }
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            return {"text": "", "error": str(e)}
    
    async def _process_pptx(self, file_path: str) -> Dict[str, Any]:
        """Extract text from PPTX."""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"--- Slide {slide_num} ---"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                text_parts.append("\n".join(slide_text))
            
            return {
                "text": "\n\n".join(text_parts),
                "num_slides": len(prs.slides),
                "source_type": "pptx"
            }
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            return {"text": "", "error": str(e)}
    
    async def _process_audio(self, file_path: str) -> Dict[str, Any]:
        """Transcribe audio using Whisper."""
        try:
            from faster_whisper import WhisperModel
            
            model = WhisperModel(
                settings.whisper_model,
                device=settings.whisper_device,
                compute_type="int8"
            )
            
            segments, info = model.transcribe(file_path, beam_size=5)
            
            text_parts = []
            for segment in segments:
                text_parts.append(segment.text)
            
            return {
                "text": " ".join(text_parts),
                "duration_seconds": int(info.duration),
                "language": info.language,
                "source_type": "audio"
            }
        except Exception as e:
            logger.error(f"Audio transcription failed: {e}")
            return {"text": "", "error": str(e)}
    
    async def _process_video(self, file_path: str) -> Dict[str, Any]:
        """Extract audio from video and transcribe."""
        try:
            # Extract audio using pydub
            from pydub import AudioSegment
            
            # Load video and extract audio
            audio = AudioSegment.from_file(file_path)
            
            # Save as temp WAV
            temp_audio = file_path + ".temp.wav"
            audio.export(temp_audio, format="wav")
            
            try:
                # Transcribe the audio
                result = await self._process_audio(temp_audio)
                result["source_type"] = "video"
                return result
            finally:
                # Clean up temp file
                if os.path.exists(temp_audio):
                    os.remove(temp_audio)
                    
        except Exception as e:
            logger.error(f"Video processing failed: {e}")
            return {"text": "", "error": str(e)}
    
    async def _process_image(self, file_path: str) -> Dict[str, Any]:
        """Extract text from image using OCR."""
        try:
            import pytesseract
            from PIL import Image
            
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            
            return {
                "text": text.strip(),
                "source_type": "image"
            }
        except Exception as e:
            logger.error(f"OCR failed: {e}")
            return {"text": "", "error": str(e)}
    
    def _chunk_text(
        self,
        text: str,
        source: str,
        material_id: str = None,
        chunk_size: int = None,
        overlap: int = None
    ) -> List[Dict[str, Any]]:
        """
        Split text into chunks for embedding.
        
        Args:
            text: Full text content
            source: Source filename
            material_id: Material ID
            chunk_size: Target chunk size in characters
            overlap: Overlap between chunks
        
        Returns:
            List of chunk dictionaries
        """
        chunk_size = chunk_size or settings.chunk_size * 4  # ~4 chars per token
        overlap = overlap or settings.chunk_overlap * 4
        
        if not text:
            return []
        
        # Split by paragraphs first
        paragraphs = text.split("\n\n")
        
        chunks = []
        current_chunk = ""
        current_start = 0
        
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            
            # If adding this paragraph exceeds chunk size
            if len(current_chunk) + len(para) > chunk_size:
                if current_chunk:
                    chunks.append({
                        "text": current_chunk.strip(),
                        "source": source,
                        "material_id": material_id,
                        "start_char": current_start,
                        "end_char": current_start + len(current_chunk),
                        "chunk_id": f"{material_id}_{len(chunks)}"
                    })
                
                # Start new chunk with overlap
                if overlap > 0 and current_chunk:
                    # Include last part of previous chunk
                    current_chunk = current_chunk[-overlap:] + "\n\n" + para
                else:
                    current_chunk = para
                current_start = current_start + len(current_chunk) - len(para)
            else:
                if current_chunk:
                    current_chunk += "\n\n" + para
                else:
                    current_chunk = para
        
        # Don't forget the last chunk
        if current_chunk.strip():
            chunks.append({
                "text": current_chunk.strip(),
                "source": source,
                "material_id": material_id,
                "start_char": current_start,
                "end_char": current_start + len(current_chunk),
                "chunk_id": f"{material_id}_{len(chunks)}"
            })
        
        return chunks
